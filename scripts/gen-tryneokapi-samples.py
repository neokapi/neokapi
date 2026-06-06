#!/usr/bin/env python3
"""Generate the real downloadable source files for the docs landing
"Try Neokapi in your browser" showcase.

Produces three tiny, valid Office/text documents that contain the showcase
term "Acme" (the default search/replace demo turns Acme → Globex). Each is
emitted as base64 so it can be embedded in
packages/kapi-playground/src/samples.ts (TRY_SAMPLES), mirroring the existing
DOCX_B64 / XLSX_B64 pattern.

The faked visual showcase in the modal needs none of these — they exist only as
the *real* "Download source" / "Download result" proof, so a visitor can open
the result in PowerPoint/Excel and confirm neokapi round-trips real files.

Usage:
    python3 scripts/gen-tryneokapi-samples.py            # write files to /tmp + print base64
    python3 scripts/gen-tryneokapi-samples.py --out DIR  # also write the raw files to DIR
"""

import argparse
import base64
import io
import os

from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches


def build_xlsx() -> bytes:
    wb = Workbook()
    ws = wb.active
    ws.title = "Q3 Report"
    ws["A1"] = "Acme quarterly revenue"
    ws["A2"] = "Acme net profit"
    ws["A3"] = "Acme customer count"
    ws["B1"] = "Total revenue"
    ws["B2"] = "Net profit"
    ws["B3"] = "Active accounts"
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def build_pptx() -> bytes:
    prs = Presentation()
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    box = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(8), Inches(1.5))
    tf = box.text_frame
    tf.text = "Welcome to Acme"
    p = tf.add_paragraph()
    p.text = "Acme makes every quarter count."
    body = slide.shapes.add_textbox(Inches(0.8), Inches(2.6), Inches(8), Inches(2))
    bf = body.text_frame
    bf.text = "Sign up for Acme today"
    p2 = bf.add_paragraph()
    p2.text = "Talk to the Acme team soon"
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def build_md() -> bytes:
    md = (
        "# Welcome to Acme\n"
        "\n"
        "Acme helps teams ship faster. Pick your favorite color and get started.\n"
        "\n"
        "- Sign up for Acme today\n"
        "- Talk to the Acme team soon\n"
    )
    return md.encode("utf-8")


def main() -> None:
    ap = argparse.ArgumentParser()
    ap.add_argument("--out", default="/tmp/tryneokapi")
    args = ap.parse_args()
    os.makedirs(args.out, exist_ok=True)

    artifacts = {
        "report.xlsx": build_xlsx(),
        "deck.pptx": build_pptx(),
        "guide.md": build_md(),
    }

    for name, data in artifacts.items():
        path = os.path.join(args.out, name)
        with open(path, "wb") as f:
            f.write(data)
        b64 = base64.b64encode(data).decode("ascii")
        const = "TRY_" + name.split(".")[-1].upper() + "_B64"
        print(f"// {name} ({len(data)} bytes)")
        print(f'export const {const} =\n  "{b64}";')
        print()
    print(f"# raw files written to {args.out}")


if __name__ == "__main__":
    main()
